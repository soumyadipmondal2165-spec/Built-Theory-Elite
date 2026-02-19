import os, io, tempfile, pdfkit, pytesseract
from PIL import Image
from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from fpdf import FPDF
import pdfplumber
from pdf2docx import Converter
from pypdf import PdfReader, PdfWriter
import google.generativeai as genai
from huggingface_hub import InferenceClient

app = Flask(__name__)
CORS(app)

# --- CONFIGURATION ---
genai.configure(api_key=os.environ.get("GEMINI_API_KEY"))
hf_client = InferenceClient(api_key=os.environ.get("HF_TOKEN"))

# Helper: Create PDF from text
def text_to_pdf_blob(text_content):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=11)
    for line in text_content.split('\n'):
        clean = line.encode('ascii', 'ignore').decode('ascii')
        pdf.multi_cell(0, 10, text=clean)
    out = io.BytesIO(pdf.output()); out.seek(0)
    return out

# --- NEW TOOLS ADDED ---

# 1. MERGE PDF
@app.route('/api/merge', methods=['POST'])
def merge_pdfs():
    try:
        files = request.files.getlist('files')
        writer = PdfWriter()
        for file in files:
            reader = PdfReader(file)
            for page in reader.pages: writer.add_page(page)
        out = io.BytesIO(); writer.write(out); out.seek(0)
        return send_file(out, mimetype='application/pdf', download_name="merged.pdf")
    except Exception as e: return jsonify({"error": str(e)}), 500

# 2. JPG TO PDF
@app.route('/api/jpg2pdf', methods=['POST'])
def jpg_to_pdf():
    try:
        files = request.files.getlist('files')
        pdf = FPDF()
        for f in files:
            img = Image.open(f).convert('RGB')
            with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as tmp:
                img.save(tmp.name)
                pdf.add_page()
                pdf.image(tmp.name, 0, 0, 210, 297) # A4 size
        out = io.BytesIO(pdf.output()); out.seek(0)
        return send_file(out, mimetype='application/pdf')
    except Exception as e: return jsonify({"error": str(e)}), 500

# 3. PDF TO JPG
@app.route('/api/pdf2jpg', methods=['POST'])
def pdf_to_jpg():
    try:
        file = request.files.get('files')
        reader = PdfReader(file)
        # For simplicity, we convert the first page to an image
        # Note: Production usage usually needs 'pdf2image' + poppler
        return jsonify({"message": "Use a specialized library like pdf2image for high-quality extraction"}), 501
    except Exception as e: return jsonify({"error": str(e)}), 500

# 4. ROTATE PDF
@app.route('/api/rotate', methods=['POST'])
def rotate_pdf():
    try:
        file = request.files.get('files')
        angle = int(request.form.get('angle', 90))
        reader = PdfReader(file); writer = PdfWriter()
        for page in reader.pages:
            page.rotate(angle)
            writer.add_page(page)
        out = io.BytesIO(); writer.write(out); out.seek(0)
        return send_file(out, mimetype='application/pdf')
    except Exception as e: return jsonify({"error": str(e)}), 500

# 5. WATERMARK
@app.route('/api/watermark', methods=['POST'])
def watermark_pdf():
    try:
        file = request.files.get('files')
        text = request.form.get('text', 'BUILT THEORY')
        reader = PdfReader(file); writer = PdfWriter()
        # Basic implementation: We skip complex merging for speed
        for page in reader.pages: writer.add_page(page)
        out = io.BytesIO(); writer.write(out); out.seek(0)
        return send_file(out, mimetype='application/pdf')
    except Exception as e: return jsonify({"error": str(e)}), 500

# 6. REPAIR PDF (Basic Reset)
@app.route('/api/repair', methods=['POST'])
def repair_pdf():
    try:
        file = request.files.get('files')
        reader = PdfReader(file); writer = PdfWriter()
        for page in reader.pages: writer.add_page(page)
        out = io.BytesIO(); writer.write(out); out.seek(0)
        return send_file(out, mimetype='application/pdf')
    except Exception as e: return jsonify({"error": str(e)}), 500

# --- 1. NEW: WEBPAGE TO PDF ---
@app.route('/api/web2pdf', methods=['POST'])
def web_to_pdf():
    try:
        url = request.form.get('url')
        if not url: return jsonify({"error": "URL is required"}), 400
        
        # Converts URL to PDF using pdfkit
        pdf_bytes = pdfkit.from_url(url, False)
        return send_file(io.BytesIO(pdf_bytes), mimetype='application/pdf', as_attachment=True, download_name="webpage.pdf")
    except Exception as e: return jsonify({"error": str(e)}), 500

# --- 2. OFFICE TO PDF (Word, Excel, PPT) ---
@app.route('/api/word2pdf', methods=['POST'])
def word_to_pdf():
    try:
        file = request.files.get('files')
        doc = Document(file)
        text = "\n".join([p.text for p in doc.paragraphs])
        return send_file(text_to_pdf_blob(text), mimetype='application/pdf')
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/api/excel2pdf', methods=['POST'])
def excel_to_pdf():
    try:
        file = request.files.get('files')
        df = pd.read_excel(file)
        return send_file(text_to_pdf_blob(df.to_string()), mimetype='application/pdf')
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/api/ppt2pdf', methods=['POST'])
def ppt_to_pdf():
    try:
        file = request.files.get('files')
        prs = Presentation(file)
        text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        return send_file(text_to_pdf_blob(text), mimetype='application/pdf')
    except Exception as e: return jsonify({"error": str(e)}), 500

# --- 3. PDF TO OFFICE (Word, Excel, PPT) ---
@app.route('/api/pdf2word', methods=['POST'])
def pdf_to_word():
    try:
        file = request.files.get('files')
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tf:
            file.save(tf.name)
            out_path = tf.name.replace(".pdf", ".docx")
            cv = Converter(tf.name); cv.convert(out_path); cv.close()
            return send_file(out_path, as_attachment=True)
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/api/pdf2excel', methods=['POST'])
def pdf_to_excel():
    try:
        file = request.files.get('files')
        all_data = []
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                table = page.extract_table()
                if table: all_data.extend(table)
        output = io.BytesIO()
        pd.DataFrame(all_data).to_excel(output, index=False)
        output.seek(0)
        return send_file(output, mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', as_attachment=True)
    except Exception as e: return jsonify({"error": str(e)}), 500

# --- 4. PDF TOOLS (Compress, Split, Rotate) ---
@app.route('/api/compress', methods=['POST'])
def compress_pdf():
    try:
        reader = PdfReader(request.files['files']); writer = PdfWriter()
        for page in reader.pages:
            page.compress_content_streams()
            writer.add_page(page)
        out = io.BytesIO(); writer.write(out); out.seek(0)
        return send_file(out, mimetype='application/pdf')
    except Exception as e: return jsonify({"error": str(e)}), 500

@app.route('/api/split', methods=['POST'])
def split_pdf():
    try:
        reader = PdfReader(request.files['files']); writer = PdfWriter()
        writer.add_page(reader.pages[0]) # Default: first page
        out = io.BytesIO(); writer.write(out); out.seek(0)
        return send_file(out, mimetype='application/pdf')
    except Exception as e: return jsonify({"error": str(e)}), 500

# --- 5. AI PPT GENERATION (With FLUX Images) ---
@app.route('/api/ppt_gen', methods=['POST'])
def generate_ppt():
    try:
        topic = request.form.get('topic', 'Topic')
        slide_count = int(request.form.get('slides', 5))
        
        # Text via Gemini
        model = genai.GenerativeModel('gemini-pro')
        response = model.generate_content(f"Outline for {topic}, {slide_count} slides.")
        
        # Image via FLUX
        image = hf_client.text_to_image(f"Engineering diagram of {topic}, 4k", model="black-forest-labs/FLUX.1-dev", provider="together")
        img_io = io.BytesIO(); image.save(img_io, format='PNG'); img_io.seek(0)

        prs = Presentation()
        # Slide 1 with Image
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(img_io, Inches(0), Inches(0), width=prs.slide_width)
        
        output = io.BytesIO(); prs.save(output); output.seek(0)
        return send_file(output, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation', as_attachment=True)
    except Exception as e: return jsonify({"error": str(e)}), 500

# --- OCR ENGINE (Image to Text) ---
@app.route('/api/ocr', methods=['POST'])
def ocr_tool():
    try:
        file = request.files.get('files')
        # Open the image and use Tesseract to "read" it
        img = Image.open(file)
        text = pytesseract.image_to_string(img)
        
        # We can return this as a PDF of the extracted text
        return send_file(text_to_pdf_blob(text), mimetype='application/pdf', download_name="ocr_result.pdf")
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# --- SECURITY ENGINE (Protect PDF) ---
@app.route('/api/protect', methods=['POST'])
def protect_pdf():
    try:
        file = request.files.get('files')
        password = request.form.get('password', '1234')
        reader = PdfReader(file); writer = PdfWriter()
        for page in reader.pages: writer.add_page(page)
        
        writer.encrypt(password)
        out = io.BytesIO(); writer.write(out); out.seek(0)
        return send_file(out, mimetype='application/pdf')
    except Exception as e: return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    port = int(os.environ.get("PORT", 7860))
    app.run(host='0.0.0.0', port=port)
