import os
import io
import tempfile
from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
import pandas as pd
from pptx import Presentation
from docx import Document
from fpdf import FPDF
import pdfplumber
from pdf2docx import Converter

app = Flask(__name__)
CORS(app) # Ensures your website can communicate with this backend

def text_to_pdf_blob(text_content):
    """Helper to convert extracted text into a valid PDF binary."""
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=11)
    for line in text_content.split('\n'):
        pdf.multi_cell(0, 10, txt=line)
    output = io.BytesIO()
    pdf.output(output)
    output.seek(0)
    return output

# --- OFFICE TO PDF TOOLS ---

@app.route('/api/word2pdf', methods=['POST'])
def word_to_pdf():
    try:
        file = request.files['files']
        doc = Document(file)
        text = "\n".join([p.text for p in doc.paragraphs])
        return send_file(text_to_pdf_blob(text), mimetype='application/pdf')
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/excel2pdf', methods=['POST'])
def excel_to_pdf():
    try:
        file = request.files['files']
        df = pd.read_excel(file)
        return send_file(text_to_pdf_blob(df.to_string()), mimetype='application/pdf')
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/ppt2pdf', methods=['POST'])
def ppt_to_pdf():
    try:
        file = request.files['files']
        prs = Presentation(file)
        text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        return send_file(text_to_pdf_blob(text), mimetype='application/pdf')
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# --- PDF TO OFFICE TOOLS ---

@app.route('/api/pdf2word', methods=['POST'])
def pdf_to_word():
    try:
        file = request.files['files']
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tf:
            file.save(tf.name)
            output_path = tf.name.replace(".pdf", ".docx")
            cv = Converter(tf.name)
            cv.convert(output_path)
            cv.close()
            return send_file(output_path, as_attachment=True)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/pdf2excel', methods=['POST'])
def pdf_to_excel():
    try:
        file = request.files['files']
        all_data = []
        with pdfplumber.open(file) as pdf:
            for page in pdf.pages:
                table = page.extract_table()
                if table: all_data.extend(table)
        output = io.BytesIO()
        pd.DataFrame(all_data).to_excel(output, index=False)
        output.seek(0)
        return send_file(output, mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', as_attachment=True)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    # Render provides the PORT environment variable automatically
    port = int(os.environ.get("PORT", 5000))
    app.run(host='0.0.0.0', port=port)